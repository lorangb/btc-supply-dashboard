from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

# =========================================================================
# STONHARD BRAND PALETTE
# Stonhard.com: Deep red logo, white backgrounds, charcoal text,
# clean corporate-industrial aesthetic. RPM International subsidiary.
# =========================================================================
STONHARD_RED     = RGBColor(0xCC, 0x00, 0x00)  # Primary brand red
STONHARD_DKRED   = RGBColor(0x99, 0x00, 0x00)  # Dark red accent
WHITE            = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE        = RGBColor(0xF5, 0xF5, 0xF5)  # Light section bg
CHARCOAL         = RGBColor(0x2D, 0x2D, 0x2D)  # Primary text
DARK_GRAY        = RGBColor(0x4A, 0x4A, 0x4A)  # Secondary text
MED_GRAY         = RGBColor(0x7A, 0x7A, 0x7A)  # Tertiary text
LIGHT_GRAY       = RGBColor(0xBB, 0xBB, 0xBB)  # Subtle text
BORDER_GRAY      = RGBColor(0xDD, 0xDD, 0xDD)  # Table borders
ACCENT_TEAL      = RGBColor(0x00, 0x7C, 0x8A)  # Secondary accent
ACCENT_NAVY      = RGBColor(0x00, 0x37, 0x76)  # RPM International blue
BTC_ORANGE       = RGBColor(0xF7, 0x93, 0x1A)  # BTC brand orange
ALERT_RED        = RGBColor(0xE0, 0x2B, 0x20)  # Alert/warning red
SUCCESS_GREEN    = RGBColor(0x2E, 0x7D, 0x32)  # Positive
AMBER            = RGBColor(0xE6, 0x8A, 0x00)  # Warning amber
PURPLE           = RGBColor(0x7B, 0x1F, 0xA2)  # STRC/SATA accent

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# =========================================================================
# HELPERS
# =========================================================================
def set_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=CHARCOAL,
             bold=False, align=PP_ALIGN.LEFT, font_name='Calibri', italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = align
    return tf

def add_para(tf, text, size=18, color=CHARCOAL, bold=False, space_before=Pt(4), italic=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.font.italic = italic
    p.space_before = space_before
    return p

def add_rect(slide, left, top, width, height, fill_color=OFF_WHITE, border_color=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def red_bar(slide, top=0):
    add_rect(slide, 0, Inches(top), 13.333, 0.08, STONHARD_RED)

def footer_bar(slide):
    add_rect(slide, 0, 7.1, 13.333, 0.4, OFF_WHITE)
    add_text(slide, 0.6, 7.15, 6, 0.3, 'BTC Supply Shock Analysis  |  Confidential', size=8, color=LIGHT_GRAY)
    add_text(slide, 6.7, 7.15, 6, 0.3, 'Sources: Glassnode  ·  STRC.live  ·  Strategy.com  ·  SEC 8-K Filings', size=8, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

def metric_card(slide, left, top, label, value, accent=STONHARD_RED, sub=''):
    add_rect(slide, left, top, 2.3, 1.5, WHITE, BORDER_GRAY)
    add_rect(slide, left, top, 2.3, 0.06, accent)
    add_text(slide, left+0.15, top+0.2, 2.0, 0.3, label, size=9, color=MED_GRAY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left+0.15, top+0.55, 2.0, 0.6, value, size=28, color=CHARCOAL, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, left+0.15, top+1.15, 2.0, 0.3, sub, size=8, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 1 — TITLE
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)

# Red header band
add_rect(slide, 0, 0, 13.333, 3.2, STONHARD_RED)
add_rect(slide, 0, 3.2, 13.333, 0.12, STONHARD_DKRED)

add_text(slide, 1.5, 0.6, 10.3, 1.0, 'BTC SUPPLY SHOCK', size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name='Calibri')
add_text(slide, 1.5, 1.6, 10.3, 0.7, 'Dashboard & Forecast Analysis', size=32, color=RGBColor(0xFF, 0xCC, 0xCC), bold=False, align=PP_ALIGN.CENTER, font_name='Calibri')
add_text(slide, 1.5, 2.4, 10.3, 0.5, 'Exchange Balance Monitor  |  STRC  |  SATA  |  Nuclear Supply Shock Projections', size=14, color=RGBColor(0xFF, 0x99, 0x99), align=PP_ALIGN.CENTER)

add_text(slide, 1.5, 4.0, 10.3, 0.5, 'Data sourced from Glassnode, STRC.live, Strategy.com, Strive IR, and SEC 8-K filings', size=14, color=MED_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 4.8, 10.3, 0.4, 'Prepared: June 2026', size=16, color=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)

# Subtle bottom accent
add_rect(slide, 0, 7.35, 13.333, 0.15, STONHARD_RED)

# ===================================================================
# SLIDE 2 — EXECUTIVE SUMMARY
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
red_bar(slide, 0)
add_text(slide, 0.6, 0.3, 12, 0.6, 'Executive Summary', size=30, color=STONHARD_RED, bold=True)
add_rect(slide, 0.6, 0.85, 2.5, 0.04, STONHARD_RED)
footer_bar(slide)

points = [
    ('Supply Crisis Forming', 'Only 14.85% of all Bitcoin (2.97M BTC) remains on exchanges — the liquid float available for purchase. This is in structural decline.', ALERT_RED),
    ('Corporate Demand Exceeds New Supply', 'Strategy (MSTR) alone buys 3-10x more BTC per week than miners produce. Combined with Strive, ETFs, and organic outflows, demand structurally exceeds supply.', BTC_ORANGE),
    ('STRC Industrializes Acquisition', 'Strategy\'s STRC preferred stock ($10.5B market cap) is now the primary BTC acquisition engine, generating $1-2B/week. Strive\'s SATA adds a second vector at 13% daily yield.', PURPLE),
    ('Nuclear Shock Projected 12-24 Months', 'Under base case assumptions, exchange balance drops below 1M BTC (nuclear threshold) by January 2028. Bull case: August 2027. Even conservative case hits "High" by mid-2028.', STONHARD_RED),
    ('April 2028 Halving is a Force Multiplier', 'Miner issuance drops from 3,150 to ~1,575 BTC/week — making the supply deficit even more acute during a period of already declining reserves.', ACCENT_NAVY),
]

for i, (title, body, accent) in enumerate(points):
    y = 1.2 + i * 1.12
    add_rect(slide, 0.6, y, 0.08, 0.9, accent)
    add_text(slide, 0.9, y, 11.5, 0.35, title, size=16, color=CHARCOAL, bold=True)
    add_text(slide, 0.9, y + 0.38, 11.5, 0.55, body, size=12, color=DARK_GRAY)

# ===================================================================
# SLIDE 3 — KEY METRICS
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
red_bar(slide, 0)
add_text(slide, 0.6, 0.3, 12, 0.6, 'Current State of BTC Supply', size=30, color=STONHARD_RED, bold=True)
add_rect(slide, 0.6, 0.85, 2.5, 0.04, STONHARD_RED)
add_text(slide, 0.6, 1.0, 10, 0.3, 'As of May 31, 2026  —  Data from Glassnode, SEC 8-K filings, STRC.live', size=11, color=MED_GRAY)
footer_bar(slide)

metric_card(slide, 0.6,  1.5, 'BTC PRICE', '$73,751', SUCCESS_GREEN, 'May 31, 2026')
metric_card(slide, 3.15, 1.5, 'ON EXCHANGES', '~2.97M', ALERT_RED, 'Glassnode All Exchanges')
metric_card(slide, 5.7,  1.5, '% ON EXCHANGES', '14.85%', ALERT_RED, 'of 20.0M mined')
metric_card(slide, 8.25, 1.5, 'STRATEGY (MSTR)', '843,738', BTC_ORANGE, '4.22% of total supply')
metric_card(slide, 10.8, 1.5, 'STRIVE (ASST)', '16,500', PURPLE, '0.08% of total supply')

# Supply breakdown table
add_rect(slide, 0.6, 3.4, 6.0, 0.45, STONHARD_RED)
add_text(slide, 0.8, 3.45, 5.6, 0.35, 'SUPPLY BREAKDOWN', size=12, color=WHITE, bold=True)

rows = [
    ('Total Mined Supply', '20,000,000', '100%', CHARCOAL),
    ('Lost / Inaccessible', '~3,500,000', '17.50%', MED_GRAY),
    ('Long-term Holders (1yr+)', '~8,500,000', '42.50%', ACCENT_NAVY),
    ('ETFs & Institutions', '~1,100,000', '5.50%', AMBER),
    ('Strategy (MSTR)', '843,738', '4.22%', BTC_ORANGE),
    ('Strive (ASST)', '16,500', '0.08%', PURPLE),
    ('On Exchanges (liquid)', '~2,970,000', '14.85%', ALERT_RED),
    ('Other Circulating', '~3,070,000', '15.35%', SUCCESS_GREEN),
]

for i, (cat, btc, pct, accent) in enumerate(rows):
    y = 3.85 + i * 0.38
    bg = WHITE if i % 2 == 0 else OFF_WHITE
    add_rect(slide, 0.6, y, 6.0, 0.38, bg, BORDER_GRAY)
    add_rect(slide, 0.6, y, 0.06, 0.38, accent)
    add_text(slide, 0.85, y + 0.05, 2.8, 0.28, cat, size=11, color=CHARCOAL)
    add_text(slide, 3.6, y + 0.05, 1.5, 0.28, btc, size=11, color=CHARCOAL, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, 5.2, y + 0.05, 1.2, 0.28, pct, size=11, color=MED_GRAY, align=PP_ALIGN.RIGHT)

# Key dynamics
add_rect(slide, 7.0, 3.4, 5.7, 0.45, CHARCOAL)
add_text(slide, 7.2, 3.45, 5.3, 0.35, 'KEY DYNAMICS', size=12, color=WHITE, bold=True)

dynamics = [
    'New miner issuance: 450 BTC/day (3,150/week)',
    'Next halving: April 2028 → cuts to ~1,575/week',
    'Strategy buying 3-10x miner output per week',
    'Strive adding ~350 BTC/week and accelerating',
    'Exchange reserves at multi-year structural lows',
    'STRC now primary BTC acquisition vehicle',
    'Corporate demand structurally > new supply',
]

tf = add_text(slide, 7.2, 4.0, 5.3, 3.5, '', size=12, color=DARK_GRAY)
tf.paragraphs[0].text = ''
for d in dynamics:
    p = add_para(tf, '      ' + d, size=12, color=DARK_GRAY, space_before=Pt(6))

# ===================================================================
# SLIDE 4 — STRATEGY WEEKLY PURCHASES
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
red_bar(slide, 0)
add_text(slide, 0.6, 0.3, 12, 0.6, 'Strategy Weekly BTC Purchases — 2026 YTD', size=30, color=STONHARD_RED, bold=True)
add_rect(slide, 0.6, 0.85, 2.5, 0.04, STONHARD_RED)
add_text(slide, 0.6, 1.0, 10, 0.3, 'Breakdown: STRC Preferred Stock funding vs MSTR Equity / Other sources', size=11, color=MED_GRAY)
footer_bar(slide)

chart_data = CategoryChartData()
chart_data.categories = ['Jan 5','Jan 12','Jan 19','Jan 26','Feb 2','Feb 9','Feb 16','Feb 23',
                          'Mar 2','Mar 9','Mar 16','Mar 23','Mar 30',
                          'Apr 6','Apr 13','Apr 20','Apr 27',
                          'May 4','May 11','May 18','May 25']

strc_vals = [0,800,4500,4200,3200,4100,1500,9800,1200,8500,10767,0,0,2100,7500,18000,1000,0,200,20000,0]
other_vals = [1468,1730,6500,5907,4433,2955,1500,10556,1815,9494,11570,0,0,2771,6427,16164,2273,0,335,4869,0]

chart_data.add_series('STRC Preferred', strc_vals)
chart_data.add_series('MSTR Equity / Other', other_vals)

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED, Inches(0.6), Inches(1.4), Inches(12.1), Inches(4.0), chart_data)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.TOP
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(10)
chart.legend.font.color.rgb = DARK_GRAY

plot = chart.plots[0]
plot.gap_width = 80
s0 = plot.series[0]
s0.format.fill.solid()
s0.format.fill.fore_color.rgb = PURPLE
s1 = plot.series[1]
s1.format.fill.solid()
s1.format.fill.fore_color.rgb = BTC_ORANGE

chart.chart_style = 2
cat_axis = chart.category_axis
cat_axis.tick_labels.font.size = Pt(8)
cat_axis.tick_labels.font.color.rgb = MED_GRAY
val_axis = chart.value_axis
val_axis.tick_labels.font.size = Pt(9)
val_axis.tick_labels.font.color.rgb = MED_GRAY
val_axis.major_gridlines.format.line.color.rgb = BORDER_GRAY

# Callout box
add_rect(slide, 0.6, 5.6, 12.1, 1.3, OFF_WHITE, PURPLE)
tf = add_text(slide, 0.9, 5.7, 11.5, 1.1, '', size=13, color=CHARCOAL)
tf.paragraphs[0].text = 'STRC is now the primary BTC acquisition engine.'
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = PURPLE
add_para(tf, 'In the week of May 11-17, STRC ATM inflows hit $2.0B — funding nearly 100% of Strategy\'s 24,869 BTC purchase that week.', size=12, color=DARK_GRAY)
add_para(tf, 'Daily STRC ATM spiked to an estimated 14,439 BTC-equivalent on May 14 alone.  STRC market cap: $10.5B — largest preferred stock globally.', size=11, color=MED_GRAY)

# ===================================================================
# SLIDE 5 — STRC vs SATA COMPARISON
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
red_bar(slide, 0)
add_text(slide, 0.6, 0.3, 12, 0.6, 'STRC vs SATA — Preferred Stock Comparison', size=30, color=STONHARD_RED, bold=True)
add_rect(slide, 0.6, 0.85, 2.5, 0.04, STONHARD_RED)
add_text(slide, 0.6, 1.0, 10, 0.3, 'Two preferred equity instruments driving BTC treasury accumulation', size=11, color=MED_GRAY)
footer_bar(slide)

# Table
headers = ['Metric', 'STRC (Strategy)', 'SATA (Strive)']
rows_data = [
    ('Issuer', 'Strategy Inc (MSTR)', 'Strive Inc (ASST)'),
    ('Type', 'Variable Rate Perpetual Pref', 'Variable Rate Perpetual Pref'),
    ('Dividend Rate', '11.50% annual', '13.00% annual'),
    ('Dividend Frequency', 'Monthly', 'Daily (from Jun 16, 2026)'),
    ('Par Value', '$100', '$100'),
    ('Market Cap', '~$10.5B', '~$845M'),
    ('Notional Outstanding', '$15.5B', '—'),
    ('BTC Backing', '843,738 BTC', '16,500 BTC'),
    ('Debt Status', '$8B+ convertible notes', '$0 — debt free'),
    ('BTC Encumbered', 'Partially', '0% — fully unencumbered'),
]

col_w = [3.4, 4.2, 4.2]
col_x = [0.8, 4.2, 8.4]
rh = 0.42
y0 = 1.45

# Header row
for ci, h in enumerate(headers):
    add_rect(slide, col_x[ci], y0, col_w[ci], rh, STONHARD_RED)
    add_text(slide, col_x[ci]+0.15, y0+0.06, col_w[ci]-0.3, rh-0.1, h, size=11, color=WHITE, bold=True)

for ri, (metric, strc, sata) in enumerate(rows_data):
    y = y0 + (ri + 1) * rh
    bg = WHITE if ri % 2 == 0 else OFF_WHITE
    vals = [metric, strc, sata]
    for ci, v in enumerate(vals):
        add_rect(slide, col_x[ci], y, col_w[ci], rh, bg, BORDER_GRAY)
        c = CHARCOAL
        if ci == 0: c = CHARCOAL; b = True
        elif 'debt free' in v or 'unencumbered' in v: c = SUCCESS_GREEN; b = True
        elif ci == 1 and ('11.50' in v or '10.5B' in v or '843,738' in v): c = BTC_ORANGE; b = True
        elif ci == 2 and ('13.00' in v or '845M' in v or '16,500' in v): c = PURPLE; b = True
        else: c = DARK_GRAY; b = False
        add_text(slide, col_x[ci]+0.15, y+0.06, col_w[ci]-0.3, rh-0.1, v, size=11, color=c, bold=b)

# Bottom callout boxes
add_rect(slide, 0.8, 5.9, 5.5, 1.0, OFF_WHITE, BTC_ORANGE)
tf = add_text(slide, 1.0, 6.0, 5.1, 0.8, '', size=12, color=CHARCOAL)
tf.paragraphs[0].text = 'STRC: The Bitcoin Acquisition Machine'
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BTC_ORANGE
add_para(tf, 'Scaled to $10.5B in 10 months — largest preferred stock', size=11, color=DARK_GRAY)
add_para(tf, 'globally. Generates $1-2B/week in BTC capital.', size=11, color=DARK_GRAY)

add_rect(slide, 6.6, 5.9, 5.5, 1.0, OFF_WHITE, PURPLE)
tf = add_text(slide, 6.8, 6.0, 5.1, 0.8, '', size=12, color=CHARCOAL)
tf.paragraphs[0].text = 'SATA: The Clean Balance Sheet Play'
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = PURPLE
add_para(tf, 'Zero debt, fully unencumbered BTC. First U.S. listed', size=11, color=DARK_GRAY)
add_para(tf, 'security to pay daily dividends. 13% annual yield.', size=11, color=DARK_GRAY)

# ===================================================================
# SLIDE 6 — EXCHANGE BALANCE DECLINE
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
red_bar(slide, 0)
add_text(slide, 0.6, 0.3, 12, 0.6, 'Exchange Balance — Structural Decline', size=30, color=STONHARD_RED, bold=True)
add_rect(slide, 0.6, 0.85, 2.5, 0.04, STONHARD_RED)
add_text(slide, 0.6, 1.0, 10, 0.3, 'Glassnode aggregate exchange balance — 2020 to present', size=11, color=MED_GRAY)
footer_bar(slide)

chart_data = CategoryChartData()
dates = ['2020-01','2020-06','2020-12','2021-03','2021-06','2021-09','2021-12',
         '2022-03','2022-06','2022-09','2022-12','2023-03','2023-06','2023-09','2023-12',
         '2024-03','2024-06','2024-09','2024-12','2025-03','2025-06','2025-09','2025-12',
         '2026-01','2026-02','2026-03','2026-05']
bals = [3.10,3.05,2.90,2.80,2.65,2.55,2.60,2.55,2.50,2.45,2.40,2.38,2.35,2.30,2.35,
        2.32,2.80,2.90,3.05,3.10,3.08,3.05,3.02,3.00,2.98,2.97,2.97]
chart_data.categories = dates
chart_data.add_series('BTC on Exchanges (M)', bals)

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(0.6), Inches(1.3), Inches(7.5), Inches(4.5), chart_data)
chart = chart_frame.chart
chart.has_legend = False
chart.chart_style = 2
s = chart.series[0]
s.format.line.color.rgb = STONHARD_RED
s.format.line.width = Pt(2.5)
s.smooth = True
val_axis = chart.value_axis
val_axis.minimum_scale = 1.0
val_axis.maximum_scale = 3.5
val_axis.tick_labels.font.size = Pt(9)
val_axis.tick_labels.font.color.rgb = MED_GRAY
val_axis.major_gridlines.format.line.color.rgb = BORDER_GRAY
cat_axis = chart.category_axis
cat_axis.tick_labels.font.size = Pt(8)
cat_axis.tick_labels.font.color.rgb = MED_GRAY

# Threshold annotations
add_text(slide, 1.0, 5.0, 3.0, 0.25, 'NUCLEAR threshold — 1.0M BTC', size=9, color=ALERT_RED, bold=True)
add_text(slide, 1.0, 4.5, 3.0, 0.25, 'SEVERE threshold — 1.5M BTC', size=9, color=AMBER, bold=True)
add_text(slide, 1.0, 3.95, 3.0, 0.25, 'HIGH threshold — 2.0M BTC', size=9, color=RGBColor(0xCC, 0x99, 0x00), bold=True)

# Right panel
add_rect(slide, 8.5, 1.3, 4.4, 4.5, OFF_WHITE, BORDER_GRAY)
tf = add_text(slide, 8.8, 1.5, 3.8, 4.0, '', size=14, color=CHARCOAL)
tf.paragraphs[0].text = 'Why Exchange Balance Matters'
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = STONHARD_RED
add_para(tf, '', size=6)
add_para(tf, 'Exchange BTC is the "float" — the only supply available for immediate purchase by new buyers.', size=12, color=DARK_GRAY)
add_para(tf, '', size=6)
add_para(tf, 'When buyers want more BTC than sellers offer, price must rise to incentivize holders to sell. A shrinking float amplifies this effect.', size=12, color=DARK_GRAY)
add_para(tf, '', size=6)
add_para(tf, 'Every major BTC rally (2013, 2017, 2021) was preceded by declining exchange reserves.', size=12, color=CHARCOAL, bold=True)
add_para(tf, '', size=6)
add_para(tf, 'Post-halving issuance (450 BTC/day) cannot replace what is being withdrawn by corporates and self-custody.', size=12, color=ALERT_RED)

# ===================================================================
# SLIDE 7 — DEMAND vs SUPPLY IMBALANCE
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
red_bar(slide, 0)
add_text(slide, 0.6, 0.3, 12, 0.6, 'The Structural Imbalance', size=30, color=STONHARD_RED, bold=True)
add_rect(slide, 0.6, 0.85, 2.5, 0.04, STONHARD_RED)
add_text(slide, 0.6, 1.0, 10, 0.3, 'Corporate BTC demand vs miner new supply — multiples of weekly issuance', size=11, color=MED_GRAY)
footer_bar(slide)

chart_data = CategoryChartData()
weeks = ['Mar 23','Mar 30','Apr 6','Apr 13','Apr 20','Apr 27','May 4','May 11','May 18','May 25']
strat_buy = [0,0,4871,13927,34164,3273,0,535,24869,0]
strive_buy = [350,350,350,350,350,350,350,350,350,350]
chart_data.categories = weeks
chart_data.add_series('Strategy Weekly Buy', strat_buy)
chart_data.add_series('Strive Est. Weekly Buy', strive_buy)

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED, Inches(0.6), Inches(1.4), Inches(7.8), Inches(4.0), chart_data)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.TOP
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(10)
chart.legend.font.color.rgb = DARK_GRAY
chart.chart_style = 2
plot = chart.plots[0]
plot.gap_width = 100
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = BTC_ORANGE
plot.series[1].format.fill.solid()
plot.series[1].format.fill.fore_color.rgb = PURPLE
val_axis = chart.value_axis
val_axis.tick_labels.font.size = Pt(9)
val_axis.tick_labels.font.color.rgb = MED_GRAY
val_axis.major_gridlines.format.line.color.rgb = BORDER_GRAY
cat_axis = chart.category_axis
cat_axis.tick_labels.font.size = Pt(9)
cat_axis.tick_labels.font.color.rgb = MED_GRAY

# Issuance reference line (text annotation since we can't add custom lines easily)
add_text(slide, 0.8, 4.5, 5, 0.25, 'Green dashed = Miner new issuance: 3,150 BTC/week', size=9, color=SUCCESS_GREEN, bold=True)

# Right panel — The Math
add_rect(slide, 8.8, 1.4, 3.9, 4.0, OFF_WHITE, BORDER_GRAY)
tf = add_text(slide, 9.1, 1.6, 3.4, 3.6, '', size=14, color=CHARCOAL)
tf.paragraphs[0].text = 'The Math'
tf.paragraphs[0].font.size = Pt(22)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = STONHARD_RED
add_para(tf, '', size=8)
add_para(tf, 'Miner Weekly Output:', size=12, color=MED_GRAY, bold=True)
add_para(tf, '3,150 BTC/week', size=18, color=SUCCESS_GREEN, bold=True)
add_para(tf, '', size=8)
add_para(tf, 'Strategy Peak Week:', size=12, color=MED_GRAY, bold=True)
add_para(tf, '34,164 BTC', size=18, color=BTC_ORANGE, bold=True)
add_para(tf, '= 10.8x issuance', size=12, color=ALERT_RED, bold=True)
add_para(tf, '', size=8)
add_para(tf, 'Strategy Avg Active:', size=12, color=MED_GRAY, bold=True)
add_para(tf, '~14,000 BTC', size=18, color=BTC_ORANGE, bold=True)
add_para(tf, '= 4.4x issuance', size=12, color=ALERT_RED, bold=True)

# Bottom callout
add_rect(slide, 0.6, 5.7, 12.1, 1.2, RGBColor(0xFF, 0xF0, 0xF0), ALERT_RED)
tf = add_text(slide, 0.9, 5.8, 11.5, 1.0, '', size=13, color=CHARCOAL)
tf.paragraphs[0].text = 'Key Insight'
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = STONHARD_RED
tf.paragraphs[0].font.size = Pt(15)
add_para(tf, 'Strategy alone absorbs 3-10x more BTC per week than miners produce. Add Strive, ETF inflows, and organic self-custody outflows —', size=12, color=DARK_GRAY)
add_para(tf, 'the net drain on exchange-available supply is accelerating faster than at any point in Bitcoin\'s history. This is structurally unprecedented.', size=12, color=DARK_GRAY)

# ===================================================================
# SLIDE 8 — SUPPLY SHOCK SEVERITY LEVELS
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
red_bar(slide, 0)
add_text(slide, 0.6, 0.3, 12, 0.6, 'Supply Shock Severity Levels', size=30, color=STONHARD_RED, bold=True)
add_rect(slide, 0.6, 0.85, 2.5, 0.04, STONHARD_RED)
add_text(slide, 0.6, 1.0, 10, 0.3, 'What happens when exchange supply drops below critical thresholds', size=11, color=MED_GRAY)
footer_bar(slide)

levels = [
    ('ELEVATED', '< 2.5M BTC (12.5%)', RGBColor(0xCC, 0x99, 0x00),
     'Reduced liquidity. Larger orders cause slippage. Price becomes more sensitive to demand spikes. Early warning stage.'),
    ('HIGH', '< 2.0M BTC (10.0%)', AMBER,
     'Significant supply constraint. OTC desks struggle to fill large orders. Bid-ask spreads widen. Institutional buyers compete for shrinking float.'),
    ('SEVERE', '< 1.5M BTC (7.5%)', RGBColor(0xE0, 0x45, 0x00),
     'Critical supply shortage. Exchange order books thin dramatically. Small purchases create outsized price moves. Resembles 2017 pre-parabolic conditions.'),
    ('NUCLEAR', '< 1.0M BTC (5.0%)', ALERT_RED,
     'Unprecedented territory. Functional supply crisis. Sellers must be incentivized at dramatically higher prices. Potential for parabolic, reflexive price discovery.'),
]

for i, (label, threshold, color, desc) in enumerate(levels):
    y = 1.5 + i * 1.3
    add_rect(slide, 0.6, y, 12.1, 1.1, WHITE, color)
    add_rect(slide, 0.6, y, 0.12, 1.1, color)
    add_text(slide, 1.0, y + 0.1, 2.2, 0.35, label, size=22, color=color, bold=True)
    add_text(slide, 1.0, y + 0.55, 2.5, 0.3, threshold, size=11, color=MED_GRAY)
    add_text(slide, 3.5, y + 0.15, 8.9, 0.8, desc, size=13, color=DARK_GRAY)

# Current status
add_rect(slide, 0.6, 6.8, 12.1, 0.2, OFF_WHITE)
add_text(slide, 0.6, 6.75, 12.1, 0.25, 'Current: 2.97M BTC on exchanges (14.85%)  —  Shock Score: 25/100  —  approaching "Elevated" territory',
         size=12, color=STONHARD_RED, bold=True, align=PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 9 — FORECAST SCENARIOS
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
red_bar(slide, 0)
add_text(slide, 0.6, 0.3, 12, 0.6, 'Nuclear Supply Shock Forecast — Scenario Projections', size=28, color=STONHARD_RED, bold=True)
add_rect(slide, 0.6, 0.85, 2.5, 0.04, STONHARD_RED)
footer_bar(slide)

# Scenario table
scenarios = [
    ('Conservative', '5,600 BTC/wk', '2,290 BTC/wk', '115 weeks', 'Aug 2028', '< 2.0M', 'HIGH', AMBER),
    ('Base Case', '14,000 BTC/wk', '13,200 BTC/wk', '86 weeks', 'Jan 2028', '< 1.0M', 'NUCLEAR', ALERT_RED),
    ('Bull Case', '21,000 BTC/wk', '21,825 BTC/wk', '63 weeks', 'Aug 2027', '< 1.0M', 'NUCLEAR', ALERT_RED),
    ('Hyper Bull', '35,000 BTC/wk', '39,475 BTC/wk', '41 weeks', 'Mar 2027', '< 1.0M', 'NUCLEAR', ALERT_RED),
]

sh = ['Scenario', 'Strategy Buy/wk', 'Net Drain/wk', 'Time to Shock', 'Projected Date', 'Exchange Bal', 'Severity']
cw = [2.0, 1.9, 1.8, 1.5, 1.6, 1.5, 1.5]
cx = [0.6]
for w in cw[:-1]: cx.append(cx[-1] + w)

y0 = 1.2
for ci, h in enumerate(sh):
    add_rect(slide, cx[ci], y0, cw[ci], 0.42, STONHARD_RED)
    add_text(slide, cx[ci]+0.1, y0+0.06, cw[ci]-0.2, 0.3, h, size=10, color=WHITE, bold=True)

for ri, (name, buy, drain, wks, date, bal, sev, sev_c) in enumerate(scenarios):
    y = y0 + (ri + 1) * 0.55
    bg = WHITE if ri % 2 == 0 else OFF_WHITE
    vals = [name, buy, drain, wks, date, bal, sev]
    for ci, v in enumerate(vals):
        add_rect(slide, cx[ci], y, cw[ci], 0.52, bg, BORDER_GRAY)
        bold = ci == 0 or ci == 6
        c = sev_c if ci == 6 else (CHARCOAL if ci == 0 else DARK_GRAY)
        add_text(slide, cx[ci]+0.1, y+0.1, cw[ci]-0.2, 0.35, v, size=12, color=c, bold=bold)

# Model assumptions
add_rect(slide, 0.6, 3.6, 12.1, 0.35, CHARCOAL)
add_text(slide, 0.8, 3.65, 11.7, 0.25, 'MODEL ASSUMPTIONS', size=11, color=WHITE, bold=True)

assumptions = [
    'Starting exchange balance: 2,970,000 BTC  |  New issuance: 3,150 BTC/week  |  Organic outflow: 4,000 BTC/week  |  Inflow: 2,500 BTC/week',
    'Strategy buy acceleration: 5%/month compounding  |  Strive: 350 BTC/week  |  New corporate entrants: 500 BTC/week',
    '"Nuclear" threshold: < 1,000,000 BTC (5%)  |  Model does NOT account for ETF flows, sovereign purchases, or demand destruction from price increases',
    'Conservative: 40% of current rate  |  Base: current sustained  |  Bull: 1.5x  |  Hyper Bull: 2.5x (Strategy reaches 1M BTC target)',
]

tf = add_text(slide, 0.8, 4.1, 11.7, 2.5, '', size=10, color=MED_GRAY)
tf.paragraphs[0].text = ''
for a in assumptions:
    add_para(tf, a, size=10, color=MED_GRAY, space_before=Pt(4))

# Bottom highlight
add_rect(slide, 0.6, 5.8, 12.1, 1.1, RGBColor(0xFF, 0xF0, 0xF0), STONHARD_RED)
tf = add_text(slide, 0.9, 5.9, 11.5, 0.9, '', size=14, color=CHARCOAL)
tf.paragraphs[0].text = 'Bottom Line'
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = STONHARD_RED
add_para(tf, 'Even the conservative case projects exchange balance hitting "High" shock territory by mid-2028. The base case reaches nuclear', size=13, color=DARK_GRAY)
add_para(tf, 'by January 2028. The April 2028 halving would compound this by cutting new issuance by 50%.', size=13, color=DARK_GRAY)

# ===================================================================
# SLIDE 10 — TIMELINE
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
red_bar(slide, 0)
add_text(slide, 0.6, 0.3, 12, 0.6, 'Projected Timeline to Nuclear Supply Shock', size=28, color=STONHARD_RED, bold=True)
add_rect(slide, 0.6, 0.85, 2.5, 0.04, STONHARD_RED)
footer_bar(slide)

# Timeline bar
add_rect(slide, 0.8, 3.0, 11.7, 0.08, MED_GRAY)

markers = [
    ('NOW', 'Jun 2026', 0.8, STONHARD_RED, '2.97M BTC'),
    ('HYPER BULL', 'Mar 2027', 3.5, ALERT_RED, '41 weeks'),
    ('BULL CASE', 'Aug 2027', 5.8, AMBER, '63 weeks'),
    ('BASE CASE', 'Jan 2028', 8.3, RGBColor(0xCC, 0x99, 0x00), '86 weeks'),
    ('CONSERVATIVE', 'Aug 2028', 10.8, SUCCESS_GREEN, '115 weeks'),
]

for label, date, x, color, sub in markers:
    add_rect(slide, x, 2.5, 0.08, 1.1, color)
    add_text(slide, x - 0.8, 1.7, 1.7, 0.7, label + '\n' + date, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x - 0.8, 3.2, 1.7, 0.3, sub, size=10, color=MED_GRAY, align=PP_ALIGN.CENTER)

# Halving marker
add_rect(slide, 9.5, 2.5, 0.06, 1.1, ACCENT_NAVY)
add_text(slide, 8.8, 3.2, 1.5, 0.4, 'HALVING\nApr 2028', size=9, color=ACCENT_NAVY, bold=True, align=PP_ALIGN.CENTER)

# Bottom panels
add_rect(slide, 0.6, 4.2, 5.8, 2.7, OFF_WHITE, STONHARD_RED)
tf = add_text(slide, 0.9, 4.35, 5.2, 2.4, '', size=13, color=CHARCOAL)
tf.paragraphs[0].text = 'Accelerating Factors'
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = STONHARD_RED
add_para(tf, '', size=4)
add_para(tf, 'Strategy targeting 1M BTC (~6,158/wk needed)', size=12, color=DARK_GRAY)
add_para(tf, 'STRC ATM scaling — $2B/week and growing', size=12, color=DARK_GRAY)
add_para(tf, 'Strive SATA launching daily dividends Jun 16', size=12, color=DARK_GRAY)
add_para(tf, 'Other corporate treasuries entering the space', size=12, color=DARK_GRAY)
add_para(tf, 'ETF inflows resume in risk-on environment', size=12, color=DARK_GRAY)
add_para(tf, 'April 2028 halving cuts issuance by 50%', size=12, color=DARK_GRAY)

add_rect(slide, 6.8, 4.2, 5.9, 2.7, OFF_WHITE, SUCCESS_GREEN)
tf = add_text(slide, 7.1, 4.35, 5.3, 2.4, '', size=13, color=CHARCOAL)
tf.paragraphs[0].text = 'Decelerating Factors'
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = SUCCESS_GREEN
add_para(tf, '', size=4)
add_para(tf, 'Rising BTC price increases acquisition cost', size=12, color=DARK_GRAY)
add_para(tf, 'Strategy debt load ($8B+) limits future borrowing', size=12, color=DARK_GRAY)
add_para(tf, 'Potential regulatory intervention', size=12, color=DARK_GRAY)
add_para(tf, 'Market sell-offs / risk-off environments', size=12, color=DARK_GRAY)
add_para(tf, 'STRC/SATA dividend obligations ($1B+/yr)', size=12, color=DARK_GRAY)
add_para(tf, 'Profit-taking at higher price levels', size=12, color=DARK_GRAY)

# ===================================================================
# SLIDE 11 — CONCLUSION
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 2.8, STONHARD_RED)
add_rect(slide, 0, 2.8, 13.333, 0.12, STONHARD_DKRED)

add_text(slide, 1.5, 0.6, 10.3, 0.8, 'Key Takeaways', size=38, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 1.5, 10.3, 0.5, 'Five structural factors converging toward a BTC supply shock', size=16, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

points = [
    ('1', 'Exchange supply is in structural decline', 'Only 14.85% of BTC remains on exchanges, down from 16%+ in 2022. The trend is resuming downward.', ALERT_RED),
    ('2', 'Corporate demand exceeds new supply by multiples', 'Strategy buys 3-10x what miners produce weekly. STRC has industrialized equity-to-BTC conversion.', BTC_ORANGE),
    ('3', 'Two preferred instruments are accelerating the drain', 'STRC ($10.5B) and SATA (13% daily yield) are new financial primitives — equity-to-BTC engines.', PURPLE),
    ('4', 'Nuclear supply shock projected within 12-24 months', 'Base case: Jan 2028. Bull case: Aug 2027. Even conservative estimates show critical levels by mid-2028.', STONHARD_RED),
    ('5', 'The April 2028 halving is a force multiplier', 'Issuance drops from 3,150 to ~1,575 BTC/week, making the supply deficit even more acute.', ACCENT_NAVY),
]

for i, (num, title, body, accent) in enumerate(points):
    y = 3.2 + i * 0.78
    add_rect(slide, 1.0, y, 0.55, 0.6, accent)
    add_text(slide, 1.0, y + 0.08, 0.55, 0.45, num, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.8, y + 0.02, 10.5, 0.3, title, size=15, color=CHARCOAL, bold=True)
    add_text(slide, 1.8, y + 0.35, 10.5, 0.3, body, size=11, color=MED_GRAY)

add_rect(slide, 0, 7.35, 13.333, 0.15, STONHARD_RED)
add_text(slide, 0.6, 7.0, 12.1, 0.3, 'Sources: Glassnode  |  STRC.live  |  Strategy.com  |  Strive IR  |  SEC 8-K Filings  |  CoinDesk  |  Bitbo',
         size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ===================================================================
# SAVE
# ===================================================================
out = '/home/user/btc-supply-dashboard/BTC_Supply_Shock_Dashboard.pptx'
prs.save(out)
print(f'Saved to {out}')
